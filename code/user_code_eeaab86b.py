from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for line in content:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0

# Prepare content for each slide
slide1_title = "Strong Statements on National Security"
slide1_content = [
    "PM Modi is making strong statements linking the recent Pahalgam attack to the historical partition and criticizing Pakistan's support for terrorism.",
    "Issued a stark warning to Pakistan, emphasizing decisive action if provoked, with statements like: 'Roti khao, warna meri goli to hai hi' ('Eat bread, or my bullet awaits').",
    "Stressed the need for evidence and proof regarding terrorism, referencing a recent incident: 'Iss baar camera ke samne kiya'."
]

slide2_title = "Comments on Indo-Pak Relations"
slide2_content = [
    "PM Modi remarked that Pakistan turned to 'proxy warfare' after failing to win a direct war against India.",
    "He commented on the Indus Water Treaty, stating, 'Haven't done anything yet and Pakistan is already sweating', indicating a tougher stance.",
    "Delivered strong messages to the people of Pakistan on tackling terrorism, urging action and accountability."
]

slide3_title = "Development Initiatives"
slide3_content = [
    "PM Modi inaugurated and laid the foundation stone for development projects worth over Rs. 53,400 crore in Bhuj, Gujarat.",
    "Reiterated his government's commitment to development and growth in Gujarat and across India.",
    "Visit highlights: Major focus on infrastructure and economic progress in the region."
]

# Create presentation
prs = Presentation()

add_slide(prs, slide1_title, slide1_content)
add_slide(prs, slide2_title, slide2_content)
add_slide(prs, slide3_title, slide3_content)

output_path = './workspace/Latest_Narendra_Modi_News.pptx'
prs.save(output_path)
print(f"PPT created at {output_path}")